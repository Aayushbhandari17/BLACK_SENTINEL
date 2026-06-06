from pptx import Presentation

def extract_text(file_path):
    presentation = Presentation(file_path)

    text_blocks = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text:
                    text_blocks.append(shape.text)

    return "\n".join(text_blocks)

class OfficeParser:
    def parse(self, file_path):
        return extract_text(file_path)